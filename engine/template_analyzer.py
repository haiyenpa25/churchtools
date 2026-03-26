#!/usr/bin/env python3
"""
template_analyzer.py — Reads a PPTX file and outputs its layout geometries 
and an AI-suggested auto-mapping of content types.
"""

import sys
import json
import os
from pptx import Presentation

def find_ph_idx(placeholders, ph_type_substrs):
    """Find the first placeholder whose type matches any string in ph_type_substrs"""
    for p in placeholders:
        if any(sub in p['type'] for sub in ph_type_substrs):
            return p['idx']
    return None

def analyze_template(pptx_path):
    if not os.path.exists(pptx_path):
        return {"status": "error", "message": f"File not found: {pptx_path}"}
        
    prs = Presentation(pptx_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    
    layouts = []
    
    for idx, layout in enumerate(prs.slide_master.slide_layouts):
        placeholders = []
        for shape in layout.shapes:
            if shape.is_placeholder:
                ph = shape.placeholder_format
                
                try:
                    ph_type = str(ph.type).split('(')[0].strip()
                except Exception:
                    ph_type = "UNKNOWN"
                    
                x, y, w, h = shape.left, shape.top, shape.width, shape.height
                
                # convert to percentage
                px = round((x / slide_w) * 100, 2)
                py = round((y / slide_h) * 100, 2)
                pw = round((w / slide_w) * 100, 2)
                ph_h = round((h / slide_h) * 100, 2)
                
                placeholders.append({
                    "idx": ph.idx,
                    "type": ph_type,
                    "x": px, "y": py, "w": pw, "h": ph_h
                })
                
        layouts.append({
            "index": idx,
            "name": layout.name,
            "placeholders": placeholders
        })
        
    # --- PHASE 1 REQ: HEURISTIC MAPPING ---
    suggested_mapping = {}
    
    for layout in layouts:
        idx = layout["index"]
        name = layout["name"].lower()
        phs = layout["placeholders"]
        
        ph_title = find_ph_idx(phs, ["TITLE", "CTR_TITLE"])
        ph_sub   = find_ph_idx(phs, ["SUBTITLE"])
        ph_body  = find_ph_idx(phs, ["BODY", "OBJECT"])
        
        bodies = [p['idx'] for p in phs if "BODY" in p['type'] or "OBJECT" in p['type']]
        ph_body2 = bodies[1] if len(bodies) > 1 else None
        
        # 1. Map 'title' (Title Slide)
        if "title" in name and "content" not in name and "title" not in suggested_mapping:
            suggested_mapping["title"] = {
                "layout_index": idx, 
                "placeholder_map": {"title": ph_title, "speaker": ph_sub, "date": ph_sub}
            }
            
        # 2. Map 'comparison' (Two Content)
        if len(bodies) >= 2 and "comparison" not in suggested_mapping:
            suggested_mapping["comparison"] = {
                "layout_index": idx,
                "placeholder_map": {"title": ph_title, "left": ph_body, "right": ph_body2}
            }
            
        # 3. Map 'content', 'list', 'main_point' (Title and Content)
        if ph_title is not None and ph_body is not None and len(bodies) == 1:
            if "content" not in suggested_mapping:
                suggested_mapping["content"] = {
                    "layout_index": idx, "placeholder_map": {"title": ph_title, "text": ph_body}
                }
            if "list" not in suggested_mapping:
                suggested_mapping["list"] = {
                    "layout_index": idx, "placeholder_map": {"title": ph_title, "items": ph_body}
                }
            if "main_point" not in suggested_mapping:
                suggested_mapping["main_point"] = {
                    "layout_index": idx, "placeholder_map": {"number": ph_title, "text": ph_body}
                }
                
        # 4. Map 'verse' & 'quote' (Section Header or Title Only)
        if ("section" in name or "quote" in name or "title only" in name) and "verse" not in suggested_mapping:
            suggested_mapping["verse"] = {
                "layout_index": idx, "placeholder_map": {"text": ph_title or ph_body, "ref": ph_sub or ph_body}
            }
            suggested_mapping["quote"] = {
                "layout_index": idx, "placeholder_map": {"text": ph_title or ph_body, "author": ph_sub or ph_body}
            }
            
    # Clean up None values in placeholder maps
    for k, mapping in suggested_mapping.items():
        mapping["placeholder_map"] = {dk: dv for dk, dv in mapping["placeholder_map"].items() if dv is not None}
        
    # Provide a fallback for blank layouts if needed
    if "blank" not in suggested_mapping:
        suggested_mapping["blank"] = {"layout_index": 6, "placeholder_map": {}}

    return {
        "status": "success",
        "slide_size": {"width_emu": slide_w, "height_emu": slide_h},
        "layouts": layouts,
        "suggested_mapping": suggested_mapping
    }

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(json.dumps({"status": "error", "message": "No input file provided."}))
        sys.exit(1)
        
    res = analyze_template(sys.argv[1])
    print(json.dumps(res))
