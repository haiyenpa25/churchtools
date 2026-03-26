"""
ppt_generator.py  —  v4 (Layout Engine Architecture)
──────────────────────────────────────────────────────
Uses layout_engine.py to resolve all coordinates from layout_schema.json.
No hard-coded inch values in this file.
"""

import sys
import json
import os
import copy

sys.path.insert(0, os.path.dirname(__file__))

from layout_engine import load_schema, resolve_layout, fit_font_size, pt_to_emu, EMU_PER_INCH

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls


# ──────────────────────────────────────────────────────────────────────────────
#  Helpers
# ──────────────────────────────────────────────────────────────────────────────

def _hex_to_rgb(hex_str: str):
    h = str(hex_str).lstrip('#').upper()
    if len(h) != 6:
        h = 'FFD700'
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def clean_text(text_val: str) -> str:
    if not text_val:
        return ""
    text_val = str(text_val).strip()
    if text_val and text_val[-1] in [',', '.', ';']:
        text_val = text_val[:-1]
    return text_val


def apply_font(run, font_def: dict, size_pt: float):
    """Apply font properties to a text run."""
    f = run.font
    f.name  = font_def.get('family', 'Georgia')
    f.size  = Pt(size_pt)
    f.bold  = True
    f.color.rgb = _hex_to_rgb(font_def.get('color_hex', 'FFD700'))

    rPr = f._element

    # Drop shadow
    if font_def.get('shadow', True):
        shadow_xml = (
            '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:outerShdw blurRad="40000" dist="35000" dir="2700000" algn="tl" rotWithShape="0">'
            '<a:srgbClr val="000000"><a:alpha val="80000"/></a:srgbClr>'
            '</a:outerShdw></a:effectLst>'
        )
        rPr.append(parse_xml(shadow_xml))

    # Stroke / outline
    if font_def.get('stroke', True):
        stroke_xml = (
            '<a:ln w="12000" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
            '</a:ln>'
        )
        rPr.append(parse_xml(stroke_xml))


def set_paragraph_spacing(paragraph, line_height: float, space_before_pt: float = 0, space_after_pt: float = 0):
    """Set line-height and paragraph spacing via XML (python-pptx lacks direct API)."""
    from pptx.oxml.ns import qn
    pPr = paragraph._p.get_or_add_pPr()

    # Line spacing: PPT uses 1/100 of % (so 112% = 112000 in hundredths or spcPct)
    # format: <a:lnSpc><a:spcPct val="112000"/></a:lnSpc>
    lnSpc = pPr.find(qn('a:lnSpc'))
    if lnSpc is not None:
        pPr.remove(lnSpc)
    lnSpc_xml = (
        f'<a:lnSpc xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:spcPct val="{round(line_height * 100000)}"/>'
        f'</a:lnSpc>'
    )
    pPr.append(parse_xml(lnSpc_xml))

    # Space before paragraph
    spc_bef = pPr.find(qn('a:spcBef'))
    if spc_bef is not None:
        pPr.remove(spc_bef)
    pPr.append(parse_xml(
        f'<a:spcBef xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:spcPts val="{round(space_before_pt * 100)}"/>'
        f'</a:spcBef>'
    ))

    # Space after paragraph
    spc_aft = pPr.find(qn('a:spcAft'))
    if spc_aft is not None:
        pPr.remove(spc_aft)
    pPr.append(parse_xml(
        f'<a:spcAft xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:spcPts val="{round(space_after_pt * 100)}"/>'
        f'</a:spcAft>'
    ))


def make_transparent_rect(slide, x: int, y: int, w: int, h: int):
    """Add a transparent rigid text container (Rectangle shape with noFill/noBorder)."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))

    spPr = shape.element.spPr

    # Remove any existing fill
    for child in list(spPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag.endswith('Fill'):
            spPr.remove(child)

    spPr.append(parse_xml(f'<a:noFill {nsdecls("a")}/>'))

    # Remove border
    existing_ln = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    if existing_ln is not None:
        spPr.remove(existing_ln)
    spPr.append(parse_xml(f'<a:ln w="0" {nsdecls("a")}><a:noFill/></a:ln>'))

    return shape


def set_txbody_margins(shape, margins: dict):
    """Set internal text body margins (inset in EMU)."""
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
    if bodyPr is None:
        return
    bodyPr.set('lIns', str(margins.get('left', 0)))
    bodyPr.set('rIns', str(margins.get('right', 0)))
    bodyPr.set('tIns', str(margins.get('top', 0)))
    bodyPr.set('bIns', str(margins.get('bottom', 0)))


# ──────────────────────────────────────────────────────────────────────────────
#  Main generator
# ──────────────────────────────────────────────────────────────────────────────

SCHEMA_PATH = os.path.join(os.path.dirname(__file__), 'layout_schema.json')


def generate_ppt(payload: dict):
    output_path = payload.get('output_file', 'output.pptx')
    template_path = payload.get('template')

    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    # Load layout schema
    schema = load_schema(SCHEMA_PATH)

    action = payload.get('action')

    if action == 'bulk_blocks':
        blocks  = payload.get('blocks', [])
        preset  = payload.get('preset', {})

        # Caller may override font_config from DB
        fc_raw = preset.get('font_config', '{}')
        if isinstance(fc_raw, str):
            try:
                fc_override = json.loads(fc_raw)
            except Exception:
                fc_override = {}
        else:
            fc_override = fc_raw or {}

        logo_override = None
        lp = preset.get('logo_path')
        if lp and os.path.exists(lp):
            logo_override = {'path': lp}

        is_green = str(preset.get('is_green_screen', '')).lower() in ['true', '1']

        for block in blocks:
            lines = block.get('lines', [])
            if not lines:
                continue

            # ─── Extract colors from fc_override (PHP already merged theme colors here) ───
            banner_color      = _hex_to_rgb(fc_override.get('banner_color',      '004D40'))
            logo_border_color = _hex_to_rgb(fc_override.get('logo_border_color', 'FFD700'))
            logo_bg_color     = _hex_to_rgb(fc_override.get('logo_bg_color',     '004D40'))
            accent_color      = _hex_to_rgb(fc_override.get('accent_color',      '002800'))
            text_color_hex    = fc_override.get('color', 'FFD700')

            # ─── Resolve layout boxes ─────────────────────────────────────────
            boxes = resolve_layout(schema, slide_w, slide_h, logo_override)

            slide_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(slide_layout)

            # ─── 1. Green screen background ───────────────────────────────────
            if is_green:
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 255, 0)

            # ─── 2. Banner rectangle ──────────────────────────────────────────
            banner_box = boxes.get('banner')
            if banner_box:
                bx = Emu(banner_box['x'])
                by = Emu(banner_box['y'])
                bw = Emu(banner_box['w'])
                bh = Emu(banner_box['h'])

                banner_shape = slide.shapes.add_shape(1, bx, by, bw, bh)  # 1 = RECTANGLE
                banner_shape.fill.solid()
                banner_shape.fill.fore_color.rgb = banner_color
                banner_shape.line.fill.background()  # no border

                # No accent stripe — bottom gap is chroma key green (transparent in livestream)

            # ─── 3. Logo circle ───────────────────────────────────────────────
            logo_box = boxes.get('logo')
            if logo_box:
                lx = Emu(logo_box['x'])
                ly = Emu(logo_box['y'])
                lw = Emu(logo_box['w'])
                lh = Emu(logo_box['h'])

                if logo_override and logo_override.get('path') and os.path.exists(logo_override['path']):
                    # User uploaded custom logo — draw as picture
                    try:
                        slide.shapes.add_picture(logo_override['path'], lx, ly, lw, lh)
                    except Exception:
                        pass
                else:
                    # Draw styled circle: background fill + gold border
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    import math
                    # Use oval (9 = MSO_SHAPE_TYPE.OVAL equivalent auto-shape index)
                    logo_shape = slide.shapes.add_shape(9, lx, ly, lw, lh)  # 9 = OVAL
                    logo_shape.fill.solid()
                    logo_shape.fill.fore_color.rgb = logo_bg_color
                    # Gold border — thickness proportional to slide
                    border_emu = int(slide_w * 0.004)
                    logo_shape.line.width = Emu(border_emu)
                    logo_shape.line.color.rgb = logo_border_color

                    # Music note symbol inside circle (centered text)
                    tf = logo_shape.text_frame
                    tf.word_wrap = False
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    para = tf.paragraphs[0]
                    para.alignment = PP_ALIGN.CENTER
                    run = para.add_run()
                    run.text = '\u266b'
                    run.font.size = Pt(22)
                    run.font.color.rgb = logo_border_color
                    run.font.bold = True

            # ─── 4. Text box ──────────────────────────────────────────────────
            txt_box = boxes.get('banner_text')
            if txt_box is None:
                continue

            font_def = txt_box['font'].copy()

            # Apply DB/override font settings
            if fc_override.get('name'):
                font_def['family'] = fc_override['name']
            font_def['color_hex'] = text_color_hex

            # Fit font size via Pillow (binary search)
            fitted_pt = fit_font_size(
                lines=lines,
                font_family=font_def['family'],
                preferred_pt=fc_override.get('size', font_def['preferred_pt']),
                min_pt=font_def['min_pt'],
                max_pt=font_def['max_pt'],
                box_w_emu=txt_box['w'],
                box_h_emu=txt_box['h'],
                line_height=font_def['line_height'],
                internal_margins_emu=txt_box['internal_margins']
            )

            shape = make_transparent_rect(
                slide,
                txt_box['x'], txt_box['y'],
                txt_box['w'], txt_box['h']
            )

            tf = shape.text_frame
            tf.word_wrap = True
            # NOTE: do NOT set auto_size — it overrides vertical_anchor and breaks MIDDLE centering

            anchor_map = {'middle': MSO_ANCHOR.MIDDLE, 'top': MSO_ANCHOR.TOP, 'bottom': MSO_ANCHOR.BOTTOM}
            tf.vertical_anchor = anchor_map.get(txt_box.get('vertical_align', 'middle'), MSO_ANCHOR.MIDDLE)

            set_txbody_margins(shape, txt_box['internal_margins'])
            tf.clear()

            align_map = {'center': PP_ALIGN.CENTER, 'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT}
            al = align_map.get(txt_box.get('text_align', 'center'), PP_ALIGN.CENTER)

            for line_obj in lines:
                p = tf.add_paragraph()
                p.alignment = al
                set_paragraph_spacing(p,
                    line_height=font_def['line_height'],
                    space_before_pt=font_def['para_space_before'],
                    space_after_pt=font_def['para_space_after'])
                primary = clean_text(line_obj.get('primary', ''))
                if primary:
                    run = p.add_run()
                    run.text = primary
                    apply_font(run, font_def, fitted_pt)

    # Save
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    print(json.dumps({"status": "success", "file": output_path}))


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(json.dumps({"status": "error", "message": "No payload provided"}))
        sys.exit(1)

    try:
        if sys.argv[1] == '--file' and len(sys.argv) >= 3:
            with open(sys.argv[2], 'r', encoding='utf-8') as f:
                payload = json.load(f)
        else:
            payload = json.loads(sys.argv[1])
        generate_ppt(payload)
    except Exception as e:
        import traceback
        print(json.dumps({"status": "error", "message": str(e), "trace": traceback.format_exc()}))
        sys.exit(1)
