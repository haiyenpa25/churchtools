import sys
import os
import json
from pptx import Presentation

def extract_text(pptx_path):
    if not os.path.exists(pptx_path):
        return {"status": "error", "message": f"Python thông báo: File vật lý hoàn toàn không tồn tại ở đường dẫn '{pptx_path}'. Có thể do lỗi lưu cache."}
        
    try:
        prs = Presentation(pptx_path)
        slides_text = []
        for slide in prs.slides:
            slide_shapes_text = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        # Clean up weird vertical tabs or carriage returns inside pptx
                        text = text.replace('\v', '\n').replace('\r', '')
                        slide_shapes_text.append(text)
            if slide_shapes_text:
                slides_text.append("\n".join(slide_shapes_text))
        
        # Join slides text with double newlines
        full_text = "\n\n".join(slides_text)
        return {"status": "success", "text": full_text}
    except Exception as e:
        return {"status": "error", "message": f"Python không thể đọc file này (Có thể do file bị hỏng, bị đặt pass, hoặc là dạng .ppt cũ bị đổi đuôi). Lỗi gốc: {str(e)}"}

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"status": "error", "message": "Missing pptx path parameter"}))
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    result = extract_text(pptx_path)
    print(json.dumps(result))
