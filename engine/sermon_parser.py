#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
sermon_parser.py — Intelligent Sermon Content Parser (Multi-Layer AI Engine)

Supports: PDF, PPTX, TXT input  |  25 slide type output
Usage:
  python sermon_parser.py --file sermon.pdf
  python sermon_parser.py --text "raw sermon text here"
  python sermon_parser.py --file payload.json   (JSON with {"text": "..."})
"""
import sys, os, json, re, argparse
os.environ.setdefault('PYTHONUTF8', '1')
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')


# ══════════════════════════════════════════════════════════════════════
#  STEP 1 — TEXT EXTRACTION
# ══════════════════════════════════════════════════════════════════════

def extract_from_pdf(path: str) -> str:
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text(x_tolerance=2) or ''
                if text.strip():
                    pages.append(text.strip())
        return '\n\n'.join(pages)
    except ImportError:
        raise RuntimeError("pdfplumber not installed. Run: pip install pdfplumber")


def extract_from_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides_text = []
        for slide in prs.slides:
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)
            if parts:
                slides_text.append('\n'.join(parts))
        return '\n\n'.join(slides_text)
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")


def extract_text(source: dict) -> str:
    if 'text' in source:
        return source['text'].strip()
    path = source.get('file', '')
    ext  = os.path.splitext(path)[1].lower()
    if ext == '.pdf':
        return extract_from_pdf(path)
    elif ext in ('.pptx', '.ppt'):
        return extract_from_pptx(path)
    elif ext == '.json':
        # Laravel passes a JSON payload file that may contain {"text": ...} or {"file": ...}
        with open(path, encoding='utf-8', errors='replace') as f:
            payload = json.load(f)
        if 'text' in payload:
            return payload['text'].strip()
        elif 'file' in payload:
            return extract_text({'file': payload['file']})  # recurse one level
        raise ValueError("JSON payload must contain 'text' or 'file' key")
    elif ext in ('.txt', '.md', ''):
        with open(path, encoding='utf-8', errors='replace') as f:
            return f.read().strip()
    raise ValueError(f"Unsupported file type: {ext}")



# ══════════════════════════════════════════════════════════════════════
#  STEP 2 — INTELLIGENT SEGMENTATION
# ══════════════════════════════════════════════════════════════════════

# Patterns that typically start a new semantic block
BLOCK_STARTERS = [
    r'^\s*\d+\s*[:\.]\s+\S',          # "1. Title" or "1: heading"
    r'^\s*[IVXLC]+\s*\.\s+\S',        # "I. Title", "II. Section"
    r'^\s*[A-Z]\s*\.\s+\S',           # "A. Point"
    r'^\s*(?:KẾT LUẬN|CONCLUSION)',     # Conclusion
    r'^\s*(?:LỜI CẦU NGUYỆN|CẦU NGUYỆN|PRAYER)',
    r'^\s*(?:ÁP DỤNG|APPLICATION)',
    r'^\s*(?:NGUYÊN NGỮ|ORIGIN WORD)',
    r'^\s*(?:TRÍCH DẪN|QUOTE)',
    r'^\s*(?:THÔNG BÁO|ANNOUNCEMENT)',
    r'^\s*(?:LỜI MỜI GỌI|INVITATION)',
]
BLOCK_STARTER_RE = re.compile('|'.join(BLOCK_STARTERS), re.IGNORECASE | re.MULTILINE)

VERSE_REF_RE = re.compile(
    r'^('
    r'(?:Sáng thế|Xuất|Lê|Dân|Phục truyền|Giô-suê|Các quan xét|Ru-tơ|I Sa|II Sa|I Vua|II Vua|'
    r'I Sử|II Sử|E-xơ-ra|Nê|Ê-xơ|Giô|Thi|Châm|Truyền|Nhã|Ê-sai|Giê|Ca|Ê-xê|Đa|Ô-sê|Giô-ên|A-mốt|'
    r'Áp|Giô-na|Mi|Na|Ha|Sô|A-ghê|Xa|Ma|Ma-thi|Mác|Lu|Giăng|Công|Rô|I Cô|II Cô|Ga|Ê-phê|Phi|'
    r'Cô|I Tê|II Tê|I Ti|II Ti|Tít|Phi-lê|Hê|Gia|I Phi|II Phi|I Giăng|II Giăng|III Giăng|Giu|Khải|'
    r'Gen|Exo|Lev|Num|Deu|Jos|Jdg|Rut|Sam|Kgs|Chr|Ezr|Neh|Est|Job|Psa|Pro|Ecc|Sol|Isa|Jer|Lam|'
    r'Eze|Dan|Hos|Joe|Amo|Oba|Jon|Mic|Nah|Hab|Zep|Hag|Zec|Mal|Mat|Mar|Luke|John|Acts|Rom|Cor|'
    r'Gal|Eph|Phil|Col|Thes|Tim|Tit|Phl|Heb|Jas|Pet|Rev'
    r')\s*\d+\s*:\s*\d+)',
    re.IGNORECASE
)

def segment_text(text: str) -> list:
    """Split text into meaningful segments."""
    # First: split by double blank lines
    raw_blocks = re.split(r'\n{2,}', text)
    segments = []
    for block in raw_blocks:
        block = block.strip()
        if not block:
            continue
        lines = block.split('\n')
        # Check if any line inside is a major heading that should split this block
        sub_segs = []
        cur = []
        for line in lines:
            stripped = line.strip()
            if stripped and BLOCK_STARTER_RE.match(stripped) and cur:
                sub_segs.append('\n'.join(cur))
                cur = [line]
            else:
                cur.append(line)
        if cur:
            sub_segs.append('\n'.join(cur))
        segments.extend([s.strip() for s in sub_segs if s.strip()])
    return segments


# ══════════════════════════════════════════════════════════════════════
#  STEP 3 — MULTI-LAYER CLASSIFIER
# ══════════════════════════════════════════════════════════════════════

HEBREW_RE   = re.compile(r'[\u05D0-\u05EA\uFB1D-\uFB4E]+')
GREEK_RE    = re.compile(r'[\u0391-\u03C9\u1F00-\u1FFF]+')
NUMBERED_RE = re.compile(r'^\s*([IVXLC]+\.|\d+\.|[A-Z]\.)\s+', re.MULTILINE)

KEYWORDS = {
    'verse': [
        'vì đức chúa trời', 'chúa phán', 'kinh thánh', 'lời chúa',
        'for god so loved', 'thus says the lord', 'the lord said',
        'bản hiệu đính', 'bản 2010', 'bản truyền thống'
    ],
    'main_point': [
        # Usually detected by structure, not keywords
    ],
    'origin': [
        'nguyên ngữ', 'tiếng hebrew', 'tiếng greek', 'tiếng hy lạp',
        'tiếng aramaic', 'phiên âm', 'có nghĩa là', 'nghĩa gốc',
        'origin word', 'word study', 'strong\'s',
    ],
    'definition': [
        'định nghĩa', 'là gì', 'khái niệm', 'thuật ngữ',
        'definition', 'what is', 'means', 'refers to',
    ],
    'application': [
        'áp dụng', 'tuần này', 'hôm nay hãy', 'thực hành',
        'apply', 'this week', 'practical', 'action step', 'challenge',
        'bác hãy', 'chúng ta hãy', 'bạn hãy', 'let us', 'let\'s',
    ],
    'prayer': [
        'lạy chúa', 'lạy cha', 'amen', 'cầu nguyện', 'xin chúa',
        'dear lord', 'heavenly father', 'we pray', 'in jesus name',
        'nhân danh chúa giê-su', 'nhân danh đức chúa giê-xu',
    ],
    'invitation': [
        'mời gọi', 'tiếp nhận', 'tin nhận', 'quyết định', 'bước lên',
        'invitation', 'come forward', 'accept christ', 'decide',
        'nếu bạn muốn', 'ai muốn tin', 'cúi đầu cầu nguyện',
    ],
    'conclusion': [
        'kết luận', 'tóm tắt', 'tóm lại', 'kết thúc', 'cuối cùng',
        'conclusion', 'summary', 'in summary', 'closing', 'to conclude',
        'hôm nay chúng ta đã', 'bài học hôm nay',
    ],
    'memory_verse': [
        'câu ghi nhớ', 'ghi nhớ', 'học thuộc', 'câu gốc tuần này',
        'memory verse', 'verse to memorize', 'take home verse',
    ],
    'testimony': [
        'nhân chứng', 'câu chuyện thật', 'chia sẻ', 'kinh nghiệm',
        'testimony', 'true story', 'personal story', 'i remember',
        'ngày kia', 'kể chuyện',
    ],
    'illustration': [
        'ví dụ', 'minh họa', 'câu chuyện', 'hình ảnh', 'tương tự',
        'illustration', 'example', 'story', 'imagine', 'picture this',
        'hãy tưởng tượng', 'ngày xưa',
    ],
    'context': [
        'bối cảnh', 'lịch sử', 'văn hóa', 'thời đó', 'vào thời',
        'background', 'historical', 'cultural context', 'in those days',
        'thế kỷ', 'năm trước công nguyên', 'năm sau công nguyên',
    ],
    'comparison': [
        'so sánh', 'đối chiếu', 'khác biệt', 'giống nhau', 'trong khi',
        'compare', 'contrast', 'difference', 'versus', ' vs ', 'on the other hand',
        'mặt khác', 'ngược lại',
    ],
    'reflection': [
        'suy ngẫm', 'tự hỏi', 'kiểm tra', 'nhìn lại',
        'reflect', 'self-examination', 'ask yourself', 'consider',
        'bạn có đang', 'chúng ta có đang',
    ],
    'quote': [
        'charles spurgeon', 'john piper', 'c.s. lewis', 'billy graham',
        'martin luther', 'timothy keller', 'famous', 'once said',
        'theologian', 'pastor', 'đã nói', 'có câu nói', 'trích dẫn',
    ],
    'song': [
        'bài hát', 'thánh ca', 'tôn vinh', 'thờ phượng',
        'song', 'hymn', 'worship', 'sing', 'let us sing',
    ],
    'announcement': [
        'thông báo', 'thông tin', 'lịch', 'sự kiện',
        'announcement', 'notice', 'upcoming event', 'save the date',
        'mời tham dự',
    ],
    'map': [
        'bản đồ', 'địa lý', 'vùng đất', 'hành trình', 'đi đến',
        'map', 'geography', 'region', 'journey to', 'located in',
    ],
    'timeline': [
        'dòng thời gian', 'mốc lịch sử', 'năm', 'thế kỷ',
        'timeline', 'chronology', 'year', 'century', 'bc', 'ad',
        'trước công nguyên', 'sau công nguyên',
    ],
}

def score_type(segment: str, candidate_type: str) -> float:
    """Return a 0.0 - 1.0 score for how likely this segment is the given type."""
    text_lower = segment.lower()
    kws = KEYWORDS.get(candidate_type, [])
    if not kws:
        return 0.0
    hits = sum(1 for kw in kws if kw in text_lower)
    return min(hits / max(len(kws) * 0.3, 1), 1.0)


def classify_segment(segment: str, index: int, total: int) -> tuple:
    """
    Returns (slide_type, confidence 0.0-1.0)
    Uses 3-layer approach: Regex → Structural → Contextual
    """
    text  = segment.strip()
    lines = text.split('\n')
    first = lines[0].strip()
    lower = text.lower()
    word_count = len(text.split())

    # ── LAYER A: Deterministic regex rules ──────────────────────────────

    # Title: first segment always classified as title
    if index == 0:
        return ('title', 0.93)

    # Verse: starts with Bible reference (MUST be before quote check)
    if VERSE_REF_RE.match(first):
        return ('verse', 0.96)

    # Inline verse: contains a reference somewhere
    if VERSE_REF_RE.search(text) and word_count < 60:
        return ('verse', 0.88)

    # Hebrew / Greek word study
    if HEBREW_RE.search(text) or GREEK_RE.search(text):
        if any(kw in lower for kw in KEYWORDS['origin']):
            return ('origin', 0.98)
        return ('origin', 0.85)

    # Keyword check — prayer (before quote)
    if any(kw in lower for kw in KEYWORDS['prayer']):
        return ('prayer', 0.90)

    # Memory verse
    if any(kw in lower for kw in KEYWORDS['memory_verse']):
        if VERSE_REF_RE.search(text):
            return ('memory_verse', 0.95)
        return ('memory_verse', 0.80)

    # Conclusion
    if any(kw in lower for kw in KEYWORDS['conclusion']):
        return ('conclusion', 0.90)

    # Invitation
    if any(kw in lower for kw in KEYWORDS['invitation']):
        return ('invitation', 0.88)

    # Application
    if any(kw in lower for kw in KEYWORDS['application']):
        return ('application', 0.87)

    # Definition
    if any(kw in lower for kw in KEYWORDS['definition']):
        return ('definition', 0.85)

    # Origin / word study
    if any(kw in lower for kw in KEYWORDS['origin']):
        return ('origin', 0.88)

    # Quote (starts and ends with quotation marks or contains attribution)
    if (text.startswith('"') and text.endswith('"')) or \
       (text.startswith('\u201c') and '\u201d' in text) or \
       re.search(r'(?:đã nói|once said|—\s*[A-ZẮẶẺẼẸ])', text):
        return ('quote', 0.88)

    # Bullet list
    bullet_count = len(re.findall(r'^[•\-\*]\s', text, re.MULTILINE))
    if bullet_count >= 2:
        return ('list', 0.88)

    # Timeline (row patterns: Year | Event or year followed by short text)
    timeline_hits = len(re.findall(r'\d{2,4}\s*(?:TC|SC|BC|AD|bn|cn)', text, re.IGNORECASE))
    if timeline_hits >= 2:
        return ('timeline', 0.90)

    # Comparison (two clear parallel sections with separator or vs/contrast)
    if any(kw in lower for kw in KEYWORDS['comparison']) and len(lines) >= 3:
        return ('comparison', 0.82)

    # Numbered list (I., II., 1., 2., A., B.)
    numbered_hits = len(NUMBERED_RE.findall(text))
    if numbered_hits >= 2:
        return ('list', 0.80)

    # Question slide (single SHORT question)
    if text.endswith('?') and word_count < 30:
        return ('question', 0.88)

    # Multiple questions
    question_count = text.count('?')
    if question_count >= 2:
        return ('reflection', 0.80)

    # Testimony / illustration keywords
    if any(kw in lower for kw in KEYWORDS['testimony']):
        return ('testimony', 0.80)
    if any(kw in lower for kw in KEYWORDS['illustration']):
        return ('illustration', 0.78)

    # Context / background
    if any(kw in lower for kw in KEYWORDS['context']):
        return ('context', 0.80)

    # ── LAYER B: Structural heuristics ──────────────────────────────────

    upper_ratio = sum(1 for c in text if c.isupper()) / max(len(text), 1)

    # ALL CAPS + short = major heading
    if upper_ratio > 0.65 and word_count <= 15:
        # Check if it has a numbering prefix (I., II., etc.)
        if re.match(r'^\s*([IVXLC]+\.|\d+\.)\s+', first, re.IGNORECASE):
            return ('main_point', 0.94)
        return ('section_break', 0.85)

    # Short single line (but mixed case) = main point or section
    if len(lines) == 1 and 3 <= word_count <= 12 and upper_ratio > 0.35:
        if re.match(r'^\s*([IVXLC]+\.|\d+\.)\s+', first, re.IGNORECASE):
            return ('main_point', 0.90)

    # Very long text = probably content/illustration
    if word_count > 100:
        if any(kw in lower for kw in KEYWORDS['illustration']):
            return ('illustration', 0.75)
        return ('content', 0.70)

    # ── LAYER C: Context-aware position rules ────────────────────────────

    if index == total - 1 or index == total - 2:
        # Last or second-to-last = conclusion or prayer
        if any(kw in lower for kw in KEYWORDS['prayer']):
            return ('prayer', 0.85)
        return ('conclusion', 0.78)

    # Fallback
    return ('content', 0.55)


# ══════════════════════════════════════════════════════════════════════
#  STEP 4 — PER-TYPE DATA EXTRACTOR
# ══════════════════════════════════════════════════════════════════════

def extract_verse(text: str) -> dict:
    lines = text.strip().split('\n')
    first_match = VERSE_REF_RE.match(lines[0].strip())
    if first_match:
        ref  = lines[0].strip()
        body = ' '.join(l.strip() for l in lines[1:] if l.strip())
        # Check for translation note
        trans_match = re.search(r'\(([^)]{3,30})\)\s*$', body)
        translation = trans_match.group(1) if trans_match else ''
        if translation:
            body = body[:body.rfind('(')].strip()
        return {'ref': ref, 'text': body, 'translation': translation}
    # Inline ref: "Giăng 3:16 — Vì Đức..."
    inline = re.split(r'\s*[—–-]\s*', text, 1)
    if len(inline) == 2 and VERSE_REF_RE.match(inline[0].strip()):
        return {'ref': inline[0].strip(), 'text': inline[1].strip(), 'translation': ''}
    return {'ref': '', 'text': text.strip(), 'translation': ''}


def extract_main_point(text: str) -> dict:
    # Detect numbering prefix
    m = re.match(r'^\s*([IVXLCM]+\.|\d+\.|[A-Z]\.)\s+(.*)', text.strip(), re.IGNORECASE | re.DOTALL)
    if m:
        return {'number': m.group(1), 'text': m.group(2).strip()}
    # ALL CAPS detection
    lines = text.strip().split('\n')
    return {'number': '', 'text': lines[0].strip()}


def extract_origin(text: str) -> dict:
    lines = [l.strip() for l in text.strip().split('\n') if l.strip()]
    # Find the actual word (Hebrew/Greek char cluster or quoted term)
    word   = ''
    phonetic = ''
    lang   = ''
    meaning = ''

    for i, line in enumerate(lines):
        # Hebrew or Greek characters
        hm = HEBREW_RE.search(line)
        gm = GREEK_RE.search(line)
        if hm: word = hm.group(0); lang = 'Hebrew'
        elif gm: word = gm.group(0); lang = 'Greek'

        # Phonetic: often in parentheses or after •
        ph = re.search(r'\(([a-z\-·]+)\)', line, re.IGNORECASE)
        if ph: phonetic = ph.group(1)

        # Meaning: after → or "nghĩa là" or "means"
        mea = re.search(r'(?:→|nghĩa là|means?:?)\s*(.+)', line, re.IGNORECASE)
        if mea: meaning = mea.group(1).strip()

    # If no Hebrew/Greek found, first non-keyword line is the word
    if not word:
        for line in lines:
            if not any(kw in line.lower() for kw in ['nguyên ngữ', 'origin', 'hebrew', 'greek']):
                word = line.strip().split()[0] if line.strip() else ''
                break

    meaning = meaning or ' '.join(
        l for l in lines[1:]
        if not HEBREW_RE.search(l) and not GREEK_RE.search(l) and '→' not in l
        and len(l.split()) > 2
    )

    return {'word': word, 'phonetic': phonetic, 'lang': lang or 'Hebrew/Greek', 'meaning': meaning}


def extract_list(text: str) -> dict:
    lines = [l.strip() for l in text.strip().split('\n') if l.strip()]
    title = ''
    items = []
    for i, line in enumerate(lines):
        clean = re.sub(r'^[•\-\*\+]\s*', '', line).strip()
        if i == 0 and not re.match(r'^[•\-\*\+]', line) and len(clean.split()) < 12:
            title = clean
        else:
            if clean:
                # Strip leading numbers  "1. item" → "item"
                clean = re.sub(r'^\d+[\.)] ?', '', clean)
                items.append(clean)
    return {'title': title, 'items': '\n'.join(items)}


def extract_comparison(text: str) -> dict:
    lines = [l.strip() for l in text.strip().split('\n') if l.strip()]
    title = lines[0] if lines else ''
    # Try to split remaining lines into two columns by separator keywords
    left_lines, right_lines = [], []
    side = 'left'
    for line in lines[1:]:
        if re.search(r'(?:nhưng|however|ngược lại|mặt khác|but|vs\.?|—)', line, re.IGNORECASE):
            side = 'right'
            continue
        if side == 'left':
            left_lines.append(line)
        else:
            right_lines.append(line)
    return {
        'title': title,
        'left':  '\n'.join(left_lines),
        'right': '\n'.join(right_lines)
    }


def extract_quote(text: str) -> dict:
    # Strip outer quotes
    clean = text.strip()
    for q_open, q_close in [('"', '"'), ('\u201c', '\u201d'), ('"', '"')]:
        if clean.startswith(q_open):
            clean = clean.lstrip(q_open)
            if q_close in clean:
                idx = clean.rfind(q_close)
                author_part = clean[idx + len(q_close):].strip()
                clean = clean[:idx].strip()
                author = re.sub(r'^[—–\-]\s*', '', author_part).strip()
                return {'text': clean, 'author': author}
    # Look for em-dash attribution at end
    m = re.search(r'[—–]\s*([^\n]+)$', text)
    author = m.group(1).strip() if m else ''
    body   = text[:m.start()].strip() if m else text.strip()
    return {'text': body, 'author': author}


def extract_timeline(text: str) -> dict:
    lines = [l.strip() for l in text.strip().split('\n') if l.strip()]
    title = ''
    events = []
    for i, line in enumerate(lines):
        if i == 0 and not re.search(r'\d{2,4}', line):
            title = line
            continue
        # Normalize separator: comma, |, dash, colon → |
        norm = re.sub(r'\s*[|,:—–]\s*', ' | ', line, count=1)
        events.append(norm)
    return {'title': title, 'events': '\n'.join(events)}


def extract_data(slide_type: str, text: str) -> dict:
    """Route to the correct per-type extractor."""
    lines = [l.strip() for l in text.strip().split('\n') if l.strip()]
    first = lines[0] if lines else ''
    body  = ' '.join(l for l in lines[1:] if l) if len(lines) > 1 else text.strip()

    if slide_type == 'title':
        # Title slide: first line = title, look for speaker/date
        speaker_m = re.search(r'(?:mục sư|ms\.?|rev\.?|pastor|diễn giả)\s*[:\.]?\s*(.+)', text, re.IGNORECASE)
        date_m    = re.search(r'\b(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4})\b', text)
        return {
            'title':    first,
            'speaker':  speaker_m.group(1).strip() if speaker_m else '',
            'date':     date_m.group(1) if date_m else '',
            'subtitle': body if not speaker_m else ''
        }

    elif slide_type == 'verse':       return extract_verse(text)
    elif slide_type == 'main_point':  return extract_main_point(text)
    elif slide_type == 'origin':      return extract_origin(text)
    elif slide_type == 'list':        return extract_list(text)
    elif slide_type == 'comparison':  return extract_comparison(text)
    elif slide_type == 'quote':       return extract_quote(text)
    elif slide_type == 'timeline':    return extract_timeline(text)

    elif slide_type == 'definition':
        m = re.search(r'(?:là|means?|nghĩa là)[:\s]+(.+)', text, re.IGNORECASE)
        return {
            'term':       first,
            'definition': m.group(1).strip() if m else body
        }

    elif slide_type in ('prayer', 'conclusion', 'reflection', 'question', 'content'):
        return {'text': text.strip()}

    elif slide_type == 'application':
        action_m = re.search(r'(?:hành động|action|tuần này|this week)[:\s]+(.+)', text, re.IGNORECASE)
        action   = action_m.group(1).strip() if action_m else ''
        main_text = text[:action_m.start()].strip() if action_m else text.strip()
        return {'text': main_text, 'action': action}

    elif slide_type == 'invitation':
        guide_m = re.search(r'(?:cúi đầu|bow|xin|please)[^\n]+', text, re.IGNORECASE)
        return {
            'text':  first,
            'guide': guide_m.group(0).strip() if guide_m else body
        }

    elif slide_type in ('illustration', 'testimony', 'context', 'image', 'map'):
        return {'title': first, 'text': body}

    elif slide_type == 'section_break':
        return {'label': first, 'subtitle': body}

    elif slide_type == 'memory_verse':
        d = extract_verse(text)
        return {'text': d.get('text', text), 'ref': d.get('ref', '')}

    elif slide_type == 'song':
        return {'name': text.strip()}

    elif slide_type == 'announcement':
        return {'title': first, 'text': body}

    elif slide_type == 'blank':
        return {}

    return {'text': text.strip()}


# ══════════════════════════════════════════════════════════════════════
#  STEP 5 — POST-PROCESSING
# ══════════════════════════════════════════════════════════════════════

def post_process(slides: list) -> list:
    """Ensure sermon structure makes sense."""
    if not slides:
        return slides

    # Ensure first slide is a title
    if slides[0]['type'] != 'title':
        # Prepend an empty title slot only if it doesn't look like a verse/point
        if slides[0]['confidence'] < 0.85:
            slides[0]['type'] = 'title'
            slides[0]['data'] = extract_data('title', slides[0].get('_raw', ''))

    # Remove very short segments that are probably page artifacts
    slides = [s for s in slides if len(s.get('_raw', '').split()) >= 2]

    # Ensure unique IDs
    for i, s in enumerate(slides):
        s['id'] = i + 1
        s.pop('_raw', None)  # Remove internal raw text

    return slides


# ══════════════════════════════════════════════════════════════════════
#  MAIN PIPELINE & AI INTEGRATION
# ══════════════════════════════════════════════════════════════════════

def parse_sermon_with_ai(source: dict, api_key: str) -> list:
    """Use Gemini (LLM) to perform Gamma-like semantic chunking."""
    raw_text = extract_text(source)
    if not raw_text.strip():
        return []
    
    try:
        import google.generativeai as genai
        genai.configure(api_key=api_key)
        
        system_instruction = '''You are an expert AI assistant for a church. Analyze the following sermon document and chunk it into presentation slides (like Gamma.app).
RULES:
1. Identify slide types: 'title', 'verse', 'main_point', 'illustration', 'prayer', 'conclusion', 'quote', 'list', 'content'.
2. DO NOT change or summarize Bible verses; extract them EXACTLY as they are written.
3. Condense long paragraphs into concise bullet points or shorter text (max 40 words per slide). If longer, split into multiple cards.
4. Output MUST be ONLY a valid JSON array of slide objects.
SCHEMA format example:
[
  {
    "type": "title", "data": {"title": "...", "subtitle": "...", "speaker": "..."}
  },
  {
    "type": "verse", "data": {"ref": "Giăng 3:16", "text": "Vì Đức Chúa Trời..."}
  },
  {
    "type": "main_point", "data": {"number": "I.", "text": "..."}
  },
  {
    "type": "list", "data": {"title": "...", "items": "bullet1\nbullet2"}
  },
  {
    "type": "quote", "data": {"text": "...", "author": "..."}
  },
  {
    "type": "content", "data": {"text": "..."}
  }
]'''
        model = genai.GenerativeModel(
            model_name="gemini-2.5-flash", 
            system_instruction=system_instruction
        )
        response = model.generate_content(
            f"Please creatively and intelligently chunk the following sermon into presentation slides suitable for screens:\n\n{raw_text}",
            generation_config=genai.types.GenerationConfig(
                temperature=0.2,
                response_mime_type="application/json"
            )
        )
        
        result_text = response.text.strip()
        slides = json.loads(result_text)
        
        # Post-process to ensure IDs and confidence
        for i, s in enumerate(slides):
            s['id'] = i + 1
            s['confidence'] = 0.98  # AI confidence is considered high
            
        return slides
    except Exception as e:
        # Fallback to rule-based engine if AI fails
        import sys
        print(json.dumps({"status": "error", "message": f"AI Parsing failed: {str(e)}", "fallback": True}), file=sys.stderr)
        return parse_sermon(source)


def parse_sermon(source: dict) -> list:
    raw_text   = extract_text(source)
    segments   = segment_text(raw_text)
    total      = len(segments)
    slides     = []

    for i, seg in enumerate(segments):
        slide_type, confidence = classify_segment(seg, i, total)
        data = extract_data(slide_type, seg)

        slide = {
            'id':         i + 1,
            'type':       slide_type,
            'data':       data,
            'confidence': round(confidence, 2),
            '_raw':       seg,
        }
        slides.append(slide)

    slides = post_process(slides)
    return slides


def main():
    parser = argparse.ArgumentParser(description='Intelligent Sermon Parser with AI')
    parser.add_argument('--file',  help='Path to PDF, PPTX, TXT, or JSON payload file')
    parser.add_argument('--text',  help='Raw sermon text to parse')
    parser.add_argument('--ai', action='store_true', help='Use Gamma-like AI semantic chunking (requires --api-key)')
    parser.add_argument('--api-key', help='Gemini API key for AI generation')
    args = parser.parse_args()

    try:
        if args.text:
            source = {'text': args.text}
        elif args.file:
            ext = os.path.splitext(args.file)[1].lower()
            if ext == '.json':
                with open(args.file, encoding='utf-8') as f:
                    payload = json.load(f)
                # If JSON has "text" key, parse that; else treat as file path list
                if 'text' in payload:
                    source = {'text': payload['text']}
                elif 'file' in payload:
                    source = {'file': payload['file']}
                else:
                    raise ValueError("JSON payload must have 'text' or 'file' key")
            else:
                source = {'file': args.file}
        else:
            # Try reading from stdin
            raw = sys.stdin.read().strip()
            if raw:
                source = {'text': raw}
            else:
                raise ValueError("No input provided. Use --file or --text")

        if args.ai and args.api_key:
            slides = parse_sermon_with_ai(source, args.api_key)
        else:
            slides = parse_sermon(source)
            
        print(json.dumps({
            'status': 'success',
            'total':  len(slides),
            'slides': slides
        }, ensure_ascii=False))

    except Exception as e:
        import traceback
        print(json.dumps({
            'status':  'error',
            'message': str(e),
            'trace':   traceback.format_exc()
        }, ensure_ascii=False))
        sys.exit(1)


if __name__ == '__main__':
    main()
