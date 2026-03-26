#!/usr/bin/env python3
"""
sermon_generator.py — Dual-export sermon PPTX engine.

Supports two export modes:
  - "live"  : Chroma-key (00FF00) background + lower-third banner
  - "full"  : Full-slide dark theme with centre-aligned content

Called from Laravel via:
  python sermon_generator.py --file /path/to/payload.json
  python sermon_generator.py --mode live --output /out/sermon_live.pptx < payload.json

JSON Payload:
{
  "mode": "live" | "full" | "both",
  "banner_ratio": 0.26,           (float 0.15-0.42, live mode only)
  "output_live": "storage/app/public/sermon_live.pptx",
  "output_full": "storage/app/public/sermon_full.pptx",
  "slides": [
    { "id": 1, "type": "title",      "data": { "title": "...", "speaker": "...", "date": "..." } },
    { "id": 2, "type": "verse",      "data": { "ref": "...", "text": "..." } },
    { "id": 3, "type": "main_point", "data": { "number": "I.", "text": "..." } },
    { "id": 4, "type": "origin",     "data": { "word": "...", "phonetic": "...", "meaning": "..." } },
    { "id": 5, "type": "list",       "data": { "title": "...", "items": "item1\nitem2" } },
    ...
  ]
}
"""

import sys, json, os
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Dimensions ──────────────────────────────────────────────────────────
SLIDE_W = 13.33   # inches  (16:9)
SLIDE_H = 7.50

def emu(i): return int(i * 914400)

# ── Colour Palette ──────────────────────────────────────────────────────
C_GREEN   = RGBColor(0x00, 0xFF, 0x00)  # Chroma-key
C_DARK    = RGBColor(0x08, 0x0A, 0x10)  # Near-black navy
C_BANNER  = RGBColor(0x0A, 0x0C, 0x16)  # Banner bg
C_GOLD    = RGBColor(0xD4, 0xA0, 0x17)  # Gold accent bar
C_YELLOW  = RGBColor(0xFF, 0xD7, 0x00)  # Bright yellow text
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY   = RGBColor(0xCB, 0xD5, 0xE1)
C_EMERALD = RGBColor(0x34, 0xD3, 0x99)
C_INDIGO  = RGBColor(0x81, 0x8C, 0xF8)
C_ROSE    = RGBColor(0xFB, 0x71, 0x85)
C_CYAN    = RGBColor(0x22, 0xD3, 0xEE)
C_AMBER   = RGBColor(0xFB, 0xBF, 0x24)
C_RED     = RGBColor(0xF8, 0x71, 0x71)
C_PURPLE  = RGBColor(0xC0, 0x84, 0xFC)
C_TEAL    = RGBColor(0x2D, 0xD4, 0xBF)
C_GREEN2  = RGBColor(0x4A, 0xDE, 0x80)

FULL_BG_TOP    = RGBColor(0x0F, 0x17, 0x2A)
FULL_BG_BOTTOM = RGBColor(0x02, 0x04, 0x0F)


# ── Helpers ─────────────────────────────────────────────────────────────
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(1, emu(x), emu(y), emu(w), emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def textbox(slide, x, y, w, h,
            text='', size=28, bold=False,
            color=C_WHITE, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE, font='Arial',
            italic=False, wrap=True):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.auto_size = None
    txb = tf._txBody
    txb.set('lIns', '72000')
    txb.set('rIns', '72000')
    txb.set('tIns', '36000')
    txb.set('bIns', '36000')
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font
    return tb


def lines_box(slide, x, y, w, h, lines: list,
              size=20, color=C_LGRAY, bullet='• ',
              accent_lines: dict = None, font='Arial'):
    """Render a bullet list of text lines inside a text-frame."""
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    txb = tf._txBody
    txb.set('lIns', '72000'); txb.set('rIns', '72000')
    txb.set('tIns', '0');     txb.set('bIns', '0')
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        c = (accent_lines or {}).get(i, color)
        run.text = (bullet if bullet else '') + str(line)
        run.font.size = Pt(size)
        run.font.color.rgb = c
        run.font.name = font
    return tb


# ════════════════════════════════════════════════════════════════════════
#  LIVE MODE — Lower-Third Builder
# ════════════════════════════════════════════════════════════════════════
def build_live_slide(prs, slide_cfg: dict, banner_h_in: float):
    slide = blank_slide(prs)
    W, H  = SLIDE_W, SLIDE_H
    bh    = banner_h_in       # banner height in inches
    by    = H - bh            # banner top Y
    AX    = 0.10              # accent bar width
    CX    = AX + 0.14         # text content start X
    CW    = W - CX - 0.25    # text content width

    t     = slide_cfg.get('type', 'content')
    d     = slide_cfg.get('data', {})

    # Chroma-key full-slide bg
    rect(slide, 0, 0, W, H, C_GREEN)

    # Banner bg + gold top accent
    rect(slide, 0, by, W, bh, C_BANNER)
    rect(slide, 0, by, W, 0.05, C_GOLD)
    # Left accent bar
    rect(slide, 0, by, AX, bh, C_GOLD)

    # Logo circle placeholder (right side)
    rect(slide, W - 0.7, by + bh * 0.15, 0.55, bh * 0.70, C_DARK)
    textbox(slide, W - 0.68, by + bh * 0.18, 0.51, bh * 0.64,
            text='LOGO', size=6, color=C_GOLD, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, font='Arial')

    CW2 = CW - 0.65   # adjust for logo space

    # ── Per-type rendering ────────────────────────────────────────────
    if t == 'title':
        textbox(slide, CX, by + 0.06, CW2, bh * 0.45,
                text=d.get('title', ''), size=32, bold=True,
                color=C_YELLOW, anchor=MSO_ANCHOR.BOTTOM)
        textbox(slide, CX, by + bh * 0.50, CW2, bh * 0.42,
                text=d.get('speaker', ''), size=18, bold=False,
                color=C_WHITE, anchor=MSO_ANCHOR.TOP)

    elif t == 'verse':
        ref = d.get('ref', '')
        textbox(slide, CX, by + 0.05, CW2, bh * 0.32,
                text=ref, size=20, bold=True,
                color=C_EMERALD, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.34, CW2, bh * 0.60,
                text=d.get('text', ''), size=20, bold=False,
                color=C_WHITE, italic=True, anchor=MSO_ANCHOR.TOP)

    elif t == 'main_point':
        num = d.get('number', '')
        prefix = (num + '  ') if num else ''
        textbox(slide, CX, by + bh * 0.08, CW2, bh * 0.84,
                text=prefix + d.get('text', ''), size=36, bold=True,
                color=C_YELLOW, anchor=MSO_ANCHOR.MIDDLE)

    elif t == 'content':
        if d.get('title'):
            textbox(slide, CX, by + 0.05, CW2, bh * 0.30,
                    text=d['title'], size=17, bold=True,
                    color=C_YELLOW, anchor=MSO_ANCHOR.MIDDLE)
            textbox(slide, CX, by + bh * 0.35, CW2, bh * 0.58,
                    text=d.get('text', ''), size=20, color=C_LGRAY,
                    anchor=MSO_ANCHOR.TOP)
        else:
            textbox(slide, CX, by + 0.08, CW2, bh - 0.16,
                    text=d.get('text', ''), size=22, color=C_WHITE,
                    anchor=MSO_ANCHOR.MIDDLE)

    elif t == 'origin':
        top = f"NGUYÊN NGỮ: '{d.get('word','').upper()}'"
        if d.get('phonetic'):
            top += f"  •  {d['phonetic']}"
        textbox(slide, CX, by + 0.05, CW2, bh * 0.38,
                text=top, size=20, bold=True,
                color=C_YELLOW, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.42, CW2, bh * 0.52,
                text=f"→ {d.get('meaning', '')}", size=18,
                color=C_WHITE, anchor=MSO_ANCHOR.TOP)

    elif t == 'definition':
        textbox(slide, CX, by + 0.05, CW2, bh * 0.38,
                text=d.get('term', '').upper(), size=24, bold=True,
                color=C_CYAN, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.42, CW2, bh * 0.52,
                text=d.get('definition', ''), size=18,
                color=C_WHITE, anchor=MSO_ANCHOR.TOP)

    elif t == 'list':
        title = d.get('title', '')
        items = [i.strip().lstrip('-').strip()
                 for i in d.get('items', '').split('\n') if i.strip()]
        visible = items[:3]
        if title:
            textbox(slide, CX, by + 0.05, CW2, bh * 0.32,
                    text=title, size=19, bold=True,
                    color=C_YELLOW, anchor=MSO_ANCHOR.MIDDLE)
            lines_box(slide, CX, by + bh * 0.36, CW2, bh * 0.58,
                      visible, size=17)
        else:
            lines_box(slide, CX, by + 0.07, CW2, bh - 0.15,
                      visible, size=19)

    elif t == 'comparison':
        textbox(slide, CX, by + 0.04, CW2, bh * 0.28,
                text=d.get('title', 'SO SÁNH'), size=18, bold=True,
                color=C_YELLOW, anchor=MSO_ANCHOR.MIDDLE)
        hw = CW2 / 2 - 0.08
        textbox(slide, CX, by + bh * 0.32, hw, bh * 0.62,
                text=d.get('left', ''), size=16, color=C_LGRAY)
        textbox(slide, CX + hw + 0.16, by + bh * 0.32, hw, bh * 0.62,
                text=d.get('right', ''), size=16, color=C_AMBER)

    elif t == 'quote':
        textbox(slide, CX, by + 0.04, CW2, bh * 0.68,
                text=f'"{d.get("text","")}"', size=20, italic=True,
                color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.72, CW2, bh * 0.24,
                text=f"— {d.get('author','')}" if d.get('author') else '',
                size=14, color=C_LGRAY, anchor=MSO_ANCHOR.MIDDLE)

    elif t == 'question':
        textbox(slide, CX, by + 0.06, CW2, bh - 0.14,
                text=d.get('text', ''), size=24, bold=False,
                color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)

    elif t == 'application':
        textbox(slide, CX, by + 0.05, CW2, bh * 0.30,
                text='✅  ÁP DỤNG HÔM NAY', size=15, bold=True,
                color=C_GREEN2, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.34, CW2, bh * 0.60,
                text=d.get('text', ''), size=19, color=C_WHITE,
                anchor=MSO_ANCHOR.TOP)

    elif t == 'conclusion':
        textbox(slide, CX, by + 0.05, CW2, bh * 0.30,
                text='KẾT LUẬN', size=18, bold=True,
                color=C_RED, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.34, CW2, bh * 0.60,
                text=d.get('text', ''), size=20, color=C_WHITE,
                anchor=MSO_ANCHOR.TOP)

    elif t == 'prayer':
        textbox(slide, CX, by + 0.05, CW2, bh - 0.12,
                text=d.get('text', ''), size=20, italic=True,
                color=C_INDIGO, anchor=MSO_ANCHOR.MIDDLE)

    elif t == 'invitation':
        textbox(slide, CX, by + 0.06, CW2, bh * 0.55,
                text=d.get('text', ''), size=26, bold=True,
                color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        if d.get('guide'):
            textbox(slide, CX, by + bh * 0.60, CW2, bh * 0.34,
                    text=d['guide'], size=15, italic=True,
                    color=C_ROSE, anchor=MSO_ANCHOR.MIDDLE)

    elif t == 'section_break':
        label = d.get('label', '')
        textbox(slide, CX, by + 0.05, CW2, bh * 0.55,
                text=label, size=34, bold=True,
                color=C_YELLOW, anchor=MSO_ANCHOR.MIDDLE)
        if d.get('subtitle'):
            textbox(slide, CX, by + bh * 0.60, CW2, bh * 0.34,
                    text=d['subtitle'], size=17, color=C_LGRAY,
                    anchor=MSO_ANCHOR.MIDDLE)

    elif t == 'memory_verse':
        textbox(slide, CX, by + 0.04, CW2, bh * 0.28,
                text='📝  CÂU GHI NHỚ', size=14, bold=True,
                color=C_PURPLE, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.32, CW2, bh * 0.42,
                text=d.get('text', ''), size=18, italic=True,
                color=C_WHITE, anchor=MSO_ANCHOR.TOP)
        textbox(slide, CX, by + bh * 0.74, CW2, bh * 0.22,
                text=d.get('ref', ''), size=14, bold=True,
                color=C_PURPLE, anchor=MSO_ANCHOR.MIDDLE)

    elif t == 'illustration':
        if d.get('title'):
            textbox(slide, CX, by + 0.04, CW2, bh * 0.30,
                    text=d['title'], size=18, bold=True,
                    color=C_AMBER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.34, CW2, bh * 0.60,
                text=d.get('text', ''), size=18, color=C_WHITE,
                anchor=MSO_ANCHOR.TOP)

    elif t == 'testimony':
        textbox(slide, CX, by + 0.04, CW2, bh * 0.28,
                text=d.get('name', ''), size=18, bold=True,
                color=C_AMBER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.32, CW2, bh * 0.62,
                text=d.get('text', ''), size=18, italic=True,
                color=C_WHITE, anchor=MSO_ANCHOR.TOP)

    elif t == 'reflection':
        textbox(slide, CX, by + 0.06, CW2, bh - 0.14,
                text=d.get('text', ''), size=21,
                color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)

    elif t == 'announcement':
        textbox(slide, CX, by + 0.04, CW2, bh * 0.36,
                text=d.get('title', ''), size=22, bold=True,
                color=C_YELLOW, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.40, CW2, bh * 0.55,
                text=d.get('text', ''), size=17, color=C_LGRAY,
                anchor=MSO_ANCHOR.TOP)

    elif t == 'song':
        textbox(slide, CX, by + 0.06, CW2, bh - 0.14,
                text=f"🎵  {d.get('name', '')}", size=26, bold=True,
                color=C_TEAL, anchor=MSO_ANCHOR.MIDDLE)

    elif t == 'context':
        if d.get('era'):
            textbox(slide, CX, by + 0.04, CW2 * 0.35, bh * 0.32,
                    text=d['era'], size=16, bold=True,
                    color=C_AMBER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.36, CW2, bh * 0.58,
                text=d.get('text', ''), size=18, color=C_LGRAY,
                anchor=MSO_ANCHOR.TOP)

    elif t == 'timeline':
        events = [e.strip() for e in d.get('events', '').split('\n') if e.strip()][:4]
        textbox(slide, CX, by + 0.04, CW2, bh * 0.28,
                text=d.get('title', ''), size=17, bold=True,
                color=C_YELLOW, anchor=MSO_ANCHOR.MIDDLE)
        lines_box(slide, CX, by + bh * 0.32, CW2, bh * 0.62,
                  events, size=15, bullet='▶  ')

    elif t == 'image':
        textbox(slide, CX, by + 0.05, CW2, bh * 0.50,
                text=f"[Hình Ảnh]", size=22, bold=True,
                color=C_LGRAY, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, CX, by + bh * 0.55, CW2, bh * 0.40,
                text=d.get('caption', ''), size=16,
                color=C_LGRAY, anchor=MSO_ANCHOR.TOP)

    elif t == 'blank':
        pass  # Just chroma-key, nothing added

    else:  # fallback
        body = d.get('text', d.get('title', str(d)))
        textbox(slide, CX, by + 0.08, CW2, bh - 0.16,
                text=body[:200], size=20, color=C_WHITE,
                anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════════
#  FULL MODE — Premium Gamma-style Slide Layout Engine
# ════════════════════════════════════════════════════════════════════════
def build_full_slide(prs, slide_cfg: dict):
    slide = blank_slide(prs)
    W, H  = SLIDE_W, SLIDE_H
    t     = slide_cfg.get('type', 'content')
    d     = slide_cfg.get('data', {})

    # PREMIUM DARK MODE BACKGROUND
    # Deep Charcoal background
    rect(slide, 0, 0, W, H, RGBColor(0x0F, 0x11, 0x1A))
    
    # Subtlety: add a bottom dark gradient/strip for grounding (Cinematic look)
    rect(slide, 0, H * 0.85, W, H * 0.15, RGBColor(0x05, 0x06, 0x0A))

    # Content area margins
    MX, MY = 1.0, 0.8
    CW, CH = W - 2*MX, H - 2*MY
    
    # Premium Typography Selection
    FONT_SERIF = 'Georgia'     # Elegant, authoritative (great for verses)
    FONT_SANS  = 'Segoe UI'    # Clean, modern (great for headers/titles)

    if t == 'title':
        # Bố cục 1: Big elegant title + subtitle
        textbox(slide, MX, H*0.25, CW, H*0.4,
                text=d.get('title', '').upper(), size=56, bold=True,
                color=RGBColor(0xF8, 0xFA, 0xFC), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM, font=FONT_SANS)
        
        # Gold underline separator
        rect(slide, W/2 - 1.5, H*0.68, 3.0, 0.02, C_GOLD)
        
        textbox(slide, MX, H*0.72, CW, H*0.12,
                text=d.get('subtitle', d.get('speaker', d.get('date', ''))).upper(), size=18,
                color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, font=FONT_SANS)

    elif t == 'verse':
        # Bố cục 2: Giant Watermark Quote for Bible verses
        textbox(slide, 0.5, 0.5, W-1, H-1, text='"', size=350, color=RGBColor(0x1E, 0x22, 0x33), 
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_SERIF)
        
        # Verse text (Serif, Italic)
        textbox(slide, MX + 0.5, MY + H*0.05, CW - 1.0, H*0.60,
                text=f"{d.get('text', '')}", size=42, italic=True,
                color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SERIF)
                
        # Verse Reference (Bottom Right Accent)
        textbox(slide, MX, H - MY - 0.4, CW, 0.5,
                text=d.get('ref', ''), size=26, bold=True,
                color=C_GOLD, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM, font=FONT_SANS)

    elif t == 'main_point':
        # Bố cục 3: Gamma Split Layout (1:2)
        num = d.get('number', '')
        if num:
            # Left column (Huge Number/Letter)
            textbox(slide, MX, MY, CW*0.18, CH,
                    text=num, size=90, bold=True, 
                    color=C_GOLD, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SERIF)
            # Right column (Text)
            textbox(slide, MX + CW*0.23, MY, CW*0.77, CH,
                    text=d.get('text', ''), size=48, bold=True,
                    color=C_WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SANS)
        else:
            # Fallback centered if no number provided
            textbox(slide, MX, MY, CW, CH,
                    text=d.get('text', ''), size=55, bold=True,
                    color=C_YELLOW, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SANS)

    elif t == 'quote':
        # Bố cục 4: Quote inside a thin gold border
        shape = slide.shapes.add_shape(1, emu(MX), emu(MY), emu(CW), emu(CH*0.85))
        shape.fill.background()
        shape.line.color.rgb = C_GOLD
        shape.line.width = Pt(1.5)
        
        textbox(slide, MX + 0.5, MY + 0.5, CW - 1.0, CH*0.85 - 1.0,
                text=f'"{d.get("text","")}"', size=38, italic=True,
                color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SERIF)
        
        if d.get('author'):
            textbox(slide, MX, MY + CH*0.85 + 0.15, CW, 0.5,
                    text=f"— {d['author']}", size=20, color=RGBColor(0x94, 0xA3, 0xB8), 
                    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SANS)

    elif t == 'list':
        # Clean list layout with checkmarks
        items = [i.strip().lstrip('-').strip() for i in d.get('items', '').split('\n') if i.strip()]
        if d.get('title'):
            textbox(slide, MX, MY, CW, H*0.15,
                    text=d['title'], size=32, bold=True,
                    color=C_GOLD, align=PP_ALIGN.LEFT, font=FONT_SANS)
        
        lines_box(slide, MX + 0.2, MY + H*0.18, CW - 0.2, CH - H*0.18, items, size=28, bullet='✓  ', color=C_WHITE, font=FONT_SANS)
        
    elif t == 'prayer':
        textbox(slide, MX, MY, CW, CH,
                text=d.get('text', ''), size=40, italic=True,
                color=RGBColor(0xE2, 0xE8, 0xF0), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SERIF)

    elif t == 'conclusion':
        # Red/Gold accent for conclusion
        textbox(slide, MX, MY, CW, H*0.15,
                text='KẾT LUẬN', size=22, bold=True,
                color=RGBColor(0xEF, 0x44, 0x44), align=PP_ALIGN.CENTER, font=FONT_SANS)
        textbox(slide, MX, MY + H*0.15, CW, CH - H*0.15,
                text=d.get('text', ''), size=45, italic=True,
                color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SERIF)

    elif t == 'question' or t == 'reflection':
        textbox(slide, MX, MY, CW, CH,
                text=d.get('text', ''), size=50, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SANS)

    elif t == 'origin' or t == 'definition':
        term = d.get('word', d.get('term', '')).upper()
        meaning = d.get('meaning', d.get('definition', ''))
        phonetic = d.get('phonetic', '')
        
        title_txt = f"NGUYÊN NGỮ: {term}" if t == 'origin' else term
        if phonetic: title_txt += f"  •  ({phonetic})"
        
        textbox(slide, MX, MY, CW, H*0.18,
                text=title_txt, size=32, bold=True,
                color=C_CYAN, align=PP_ALIGN.CENTER, font=FONT_SANS)
        textbox(slide, MX, MY + H*0.2, CW, CH - H*0.2,
                text=f"→ {meaning}", size=36,
                color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SERIF)

    elif t == 'image':
        rect(slide, MX, MY, CW, CH*0.8, RGBColor(0x1E, 0x29, 0x3B))
        textbox(slide, MX, MY, CW, CH*0.8,
                text="[Media Placeholder]", size=24, color=RGBColor(0x47, 0x55, 0x69), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SANS)
        if d.get('caption'):
            textbox(slide, MX, MY + CH*0.85, CW, CH*0.15,
                    text=d['caption'], size=18, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER, font=FONT_SANS)

    elif t == 'blank':
        pass

    else:
        # Fallback elegant content rendering
        title = d.get('title', d.get('name', d.get('term', '')))
        body  = d.get('text', d.get('events', ''))
        if title:
            textbox(slide, MX, MY, CW, H*0.15,
                    text=title.upper(), size=24, bold=True,
                    color=C_GOLD, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SANS)
            
        textbox(slide, MX, MY + (H*0.18 if title else 0), CW, CH - (H*0.18 if title else 0),
                text=body[:350], size=32, color=C_WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=FONT_SANS)


# ════════════════════════════════════════════════════════════════════════
#  TEMPLATE MODE — Native PowerPoint Master Template Mapping
# ════════════════════════════════════════════════════════════════════════
def build_template_slide(prs, slide_cfg: dict, template_mapping: dict):
    t = slide_cfg.get('type', 'content')
    d = slide_cfg.get('data', {})
    
    mapping = template_mapping.get(t, {})
    layout_index = mapping.get('layout_index', 6)  # Default to blank if not found
    
    try:
        slide_layout = prs.slide_layouts[layout_index]
    except IndexError:
        slide_layout = prs.slide_layouts[6]
        
    slide = prs.slides.add_slide(slide_layout)
    
    # Gather all text placeholders
    phs = [shape for shape in slide.placeholders if shape.has_text_frame]
    
    # 1. If explicit placeholder mapping is provided in the schema
    ph_map = mapping.get('placeholder_map', {})
    if ph_map:
        for data_key, ph_idx in ph_map.items():
            val = d.get(data_key)
            if val is not None:
                for p in phs:
                    if p.placeholder_format.idx == int(ph_idx):
                        p.text = str(val)
                        break
    else:
        # 2. Heuristic fallback mapping
        phs.sort(key=lambda x: x.top)
        
        # Compile content
        title_val = d.get('title', d.get('term', d.get('name', d.get('word', ''))))
        body_val = d.get('text', d.get('definition', d.get('meaning', d.get('events', ''))))
        ref_val = d.get('ref', '')
        
        # Try finding dedicated TITLE placeholders (type 1 or 3)
        title_has_been_set = False
        if title_val:
            for p in phs:
                if p.placeholder_format.type in (1, 3):
                    p.text = str(title_val)
                    title_has_been_set = True
                    break
        
        # Remaining values to fill into the rest of the boxes
        leftovers = []
        if title_val and not title_has_been_set: leftovers.append(title_val)
        if body_val: leftovers.append(body_val)
        if ref_val: leftovers.append(ref_val)
        if d.get('speaker'): leftovers.append(d.get('speaker'))
        if d.get('author'): leftovers.append(d.get('author'))
        
        for p in phs:
            if not p.text and leftovers:
                p.text = str(leftovers.pop(0))


# ════════════════════════════════════════════════════════════════════════
#  MAIN ENTRY
# ════════════════════════════════════════════════════════════════════════
def new_prs():
    prs = Presentation()
    prs.slide_width  = Emu(emu(SLIDE_W))
    prs.slide_height = Emu(emu(SLIDE_H))
    return prs


def generate(payload: dict) -> dict:
    mode        = payload.get('mode', 'both')
    banner_h    = float(payload.get('banner_ratio', 0.26)) * SLIDE_H
    slides_data = payload.get('slides', [])
    out_live    = payload.get('output_live', 'storage/app/public/sermon_live.pptx')
    out_full    = payload.get('output_full', 'storage/app/public/sermon_full.pptx')
    
    # Template Integration Payload Additions
    template_file = payload.get('template_file')
    template_mapping = payload.get('template_mapping', {})

    result = {}

    if mode in ('live', 'both'):
        prs = new_prs()
        for s in slides_data:
            build_live_slide(prs, s, banner_h)
        os.makedirs(os.path.dirname(os.path.abspath(out_live)), exist_ok=True)
        prs.save(out_live)
        result['live'] = out_live

    if mode in ('full', 'both'):
        if template_file and os.path.exists(template_file):
            prs = Presentation(template_file)
            for s in slides_data:
                build_template_slide(prs, s, template_mapping)
        else:
            prs = new_prs()
            for s in slides_data:
                build_full_slide(prs, s)
                
        os.makedirs(os.path.dirname(os.path.abspath(out_full)), exist_ok=True)
        prs.save(out_full)
        result['full'] = out_full

    return result


if __name__ == '__main__':
    try:
        if '--file' in sys.argv:
            idx = sys.argv.index('--file')
            with open(sys.argv[idx + 1], 'r', encoding='utf-8') as f:
                payload = json.load(f)
        else:
            payload = json.loads(sys.argv[1] if len(sys.argv) > 1 else sys.stdin.read())

        result = generate(payload)
        print(json.dumps({'status': 'success', **result}, ensure_ascii=False))
    except Exception as e:
        import traceback
        print(json.dumps({'status': 'error', 'message': str(e),
                          'trace': traceback.format_exc()}))
        sys.exit(1)
